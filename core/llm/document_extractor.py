"""
Document Content Extraction Utilities

Extracts text content from various document formats:
- PDF files (PyPDF2)
- Word documents - DOCX (python-docx)
- Excel spreadsheets - XLSX, XLS (openpyxl, pandas)
- PowerPoint presentations - PPTX (python-pptx)
- LibreOffice documents - ODT, ODS, ODP

Used to make document contents accessible to LLM models that don't have native
document parsing capabilities.
"""

import io
import logging
from typing import Optional, Dict, Any
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_pdf_content(file_data: bytes, filename: str) -> str:
    """
    Extract text content from PDF files using PyPDF2.

    Args:
        file_data: Binary content of the PDF file
        filename: Name of the file (for logging)

    Returns:
        Extracted text content

    Raises:
        Exception: If extraction fails
    """
    try:
        import PyPDF2

        pdf_file = io.BytesIO(file_data)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        # Extract text from all pages
        text_content = []
        for page_num, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            if page_text.strip():
                text_content.append(f"--- Page {page_num} ---\n{page_text}")

        if not text_content:
            return "⚠️ PDF file appears to be empty or contains no extractable text."

        extracted_text = "\n\n".join(text_content)
        logger.info(f"Successfully extracted {len(extracted_text)} characters from PDF: {filename}")
        return extracted_text

    except Exception as e:
        logger.error(f"Failed to extract PDF content from {filename}: {str(e)}")
        raise Exception(f"Could not extract PDF content: {str(e)}")


def extract_docx_content(file_data: bytes, filename: str) -> str:
    """
    Extract text content from DOCX files using python-docx.

    Args:
        file_data: Binary content of the DOCX file
        filename: Name of the file (for logging)

    Returns:
        Extracted text content

    Raises:
        Exception: If extraction fails
    """
    try:
        from docx import Document

        docx_file = io.BytesIO(file_data)
        doc = Document(docx_file)

        # Extract text from all paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract text from tables
        tables_text = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(row_data):
                    table_data.append(" | ".join(row_data))
            if table_data:
                tables_text.append("\n".join(table_data))

        content_parts = []
        if paragraphs:
            content_parts.append("\n\n".join(paragraphs))
        if tables_text:
            content_parts.append("\n\n--- Tables ---\n\n" + "\n\n".join(tables_text))

        if not content_parts:
            return "⚠️ Word document appears to be empty."

        extracted_text = "\n\n".join(content_parts)
        logger.info(f"Successfully extracted {len(extracted_text)} characters from DOCX: {filename}")
        return extracted_text

    except Exception as e:
        logger.error(f"Failed to extract DOCX content from {filename}: {str(e)}")
        raise Exception(f"Could not extract Word document content: {str(e)}")


def extract_xlsx_content(file_data: bytes, filename: str) -> str:
    """
    Extract text content from XLSX files using openpyxl and pandas.

    Args:
        file_data: Binary content of the XLSX file
        filename: Name of the file (for logging)

    Returns:
        Extracted text content with sheet names and data

    Raises:
        Exception: If extraction fails
    """
    try:
        import openpyxl  # type: ignore[import-untyped]
        import pandas as pd  # type: ignore[import-untyped]

        excel_file = io.BytesIO(file_data)
        workbook = openpyxl.load_workbook(excel_file, data_only=True)

        sheets_content = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # Get data from sheet
            data = []
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                if any(cell is not None and str(cell).strip() for cell in row):
                    data.append([str(cell) if cell is not None else "" for cell in row])

            if data:
                # Convert to DataFrame for better formatting
                df = pd.DataFrame(data[1:], columns=data[0] if data else None)
                sheets_content.append(f"--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}")

        if not sheets_content:
            return "⚠️ Excel file appears to be empty."

        extracted_text = "\n\n".join(sheets_content)
        logger.info(f"Successfully extracted {len(extracted_text)} characters from XLSX: {filename}")
        return extracted_text

    except Exception as e:
        logger.error(f"Failed to extract XLSX content from {filename}: {str(e)}")
        raise Exception(f"Could not extract Excel content: {str(e)}")


def extract_pptx_content(file_data: bytes, filename: str) -> str:
    """
    Extract text content from PPTX files using python-pptx.

    Args:
        file_data: Binary content of the PPTX file
        filename: Name of the file (for logging)

    Returns:
        Extracted text content with slide numbers

    Raises:
        Exception: If extraction fails
    """
    try:
        from pptx import Presentation

        pptx_file = io.BytesIO(file_data)
        prs = Presentation(pptx_file)

        slides_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                slides_content.append(f"--- Slide {slide_num} ---\n" + "\n\n".join(slide_text))

        if not slides_content:
            return "⚠️ PowerPoint presentation appears to be empty or contains no text."

        extracted_text = "\n\n".join(slides_content)
        logger.info(f"Successfully extracted {len(extracted_text)} characters from PPTX: {filename}")
        return extracted_text

    except Exception as e:
        logger.error(f"Failed to extract PPTX content from {filename}: {str(e)}")
        raise Exception(f"Could not extract PowerPoint content: {str(e)}")


def extract_document_content(file_data: bytes, filename: str, mime_type: Optional[str] = None) -> Dict[str, Any]:
    """
    Extract content from a document file based on its extension or MIME type.

    Args:
        file_data: Binary content of the file
        filename: Name of the file
        mime_type: Optional MIME type of the file

    Returns:
        Dictionary with:
            - success: bool
            - content: str (extracted text content)
            - error: Optional[str] (error message if failed)
            - file_type: str (detected file type)
    """
    # Determine file type from extension
    ext = Path(filename).suffix.lower()

    try:
        # PDF files
        if ext == '.pdf' or mime_type == 'application/pdf':
            content = extract_pdf_content(file_data, filename)
            return {
                'success': True,
                'content': content,
                'error': None,
                'file_type': 'pdf'
            }

        # Word documents
        elif ext in ['.docx', '.docm'] or mime_type in [
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.ms-word.document.macroEnabled.12'
        ]:
            content = extract_docx_content(file_data, filename)
            return {
                'success': True,
                'content': content,
                'error': None,
                'file_type': 'docx'
            }

        # Excel spreadsheets
        elif ext in ['.xlsx', '.xlsm', '.xlsb'] or mime_type in [
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-excel.sheet.macroEnabled.12',
            'application/vnd.ms-excel.sheet.binary.macroEnabled.12'
        ]:
            content = extract_xlsx_content(file_data, filename)
            return {
                'success': True,
                'content': content,
                'error': None,
                'file_type': 'xlsx'
            }

        # PowerPoint presentations
        elif ext in ['.pptx', '.pptm'] or mime_type in [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint.presentation.macroEnabled.12'
        ]:
            content = extract_pptx_content(file_data, filename)
            return {
                'success': True,
                'content': content,
                'error': None,
                'file_type': 'pptx'
            }

        # Unsupported file type
        else:
            error_msg = f"Unsupported file type: {ext}. Supported types: PDF, DOCX, XLSX, PPTX"
            logger.warning(f"{error_msg} for file: {filename}")
            return {
                'success': False,
                'content': None,
                'error': error_msg,
                'file_type': 'unknown'
            }

    except Exception as e:
        error_msg = str(e)
        logger.error(f"Document extraction failed for {filename}: {error_msg}")
        return {
            'success': False,
            'content': None,
            'error': error_msg,
            'file_type': ext[1:] if ext else 'unknown'
        }
